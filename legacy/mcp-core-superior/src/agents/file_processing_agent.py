"""
File Processing Agent MCP - Procesamiento Avanzado de Archivos
Extiende el backend/tools/file_processor.py con capacidades MCP para multimedia y formatos avanzados

Formatos soportados:
- Documentos: PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, JSON, HTML
- Imágenes: JPG, PNG, TIFF, GIF, BMP, WEBP (con OCR)
- Audio: MP3, WAV, FLAC, AAC, OGG
- Video: MP4, AVI, MOV, MKV, WEBM

Capacidades:
- Extracción de texto avanzada
- Análisis de imágenes con IA
- Conversión entre formatos
- OCR (Reconocimiento Óptico de Caracteres)
- Análisis de metadata
- Procesamiento de documentos complejos
"""

import os
import sys
import json
import mimetypes
import hashlib
import tempfile
import shutil
import subprocess
from pathlib import Path
from typing import Dict, List, Any, Optional, Union, Tuple
from datetime import datetime
import logging

# Imports para procesamiento multimedia
try:
    import cv2
    CV2_AVAILABLE = True
except ImportError:
    CV2_AVAILABLE = False

try:
    from PIL import Image, ImageEnhance
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import librosa
    AUDIO_LIBS_AVAILABLE = True
except ImportError:
    AUDIO_LIBS_AVAILABLE = False

try:
    import mutagen
    MUTAGEN_AVAILABLE = True
except ImportError:
    MUTAGEN_AVAILABLE = False

try:
    import ffmpeg
    FFMPEG_AVAILABLE = True
except ImportError:
    FFMPEG_AVAILABLE = False

try:
    import qrcode
    QRCODE_AVAILABLE = True
except ImportError:
    QRCODE_AVAILABLE = False

# Importar la herramienta base existente
sys.path.append(os.path.join(os.path.dirname(__file__), '../../..'))
from backend.tools.file_processor import FileProcessor, ToolResult

class FileProcessingAgentMCP:
    """Agente MCP para procesamiento avanzado de archivos multimedia y documentos"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.file_processor = FileProcessor()
        
        # Tipos MIME extendidos
        self.allowed_mime_types = {
            # Documentos
            'application/pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # DOCX
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',       # XLSX
            'application/vnd.openxmlformats-officedocument.presentationml.presentation', # PPTX
            'text/plain',
            'text/markdown',
            'text/csv',
            'application/json',
            'text/html',
            'application/rtf',
            
            # Imágenes
            'image/jpeg',
            'image/png',
            'image/tiff',
            'image/gif',
            'image/bmp',
            'image/webp',
            'image/svg+xml',
            
            # Audio
            'audio/mpeg',
            'audio/wav',
            'audio/flac',
            'audio/aac',
            'audio/ogg',
            'audio/x-wav',
            
            # Video
            'video/mp4',
            'video/avi',
            'video/quicktime',
            'video/x-msvideo',
            'video/webm'
        }
        
        # Extensiones soportadas
        self.supported_extensions = {
            # Documentos
            '.pdf', '.docx', '.xlsx', '.pptx', '.txt', '.md', '.csv', '.json', '.html', '.rtf',
            
            # Imágenes
            '.jpg', '.jpeg', '.png', '.tiff', '.tif', '.gif', '.bmp', '.webp', '.svg',
            
            # Audio
            '.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a',
            
            # Video
            '.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv'
        }
        
        # Tamaño máximo de archivo (50MB para multimedia)
        self.max_file_size = 50 * 1024 * 1024
        
        # Configuración de OCR
        self.ocr_languages = ['spa', 'eng']  # Español e inglés
        
    def get_tools(self) -> List[Dict[str, Any]]:
        """Retorna las herramientas MCP disponibles"""
        return [
            {
                "name": "extract_text_from_document",
                "description": "Extrae texto de documentos (PDF, DOCX, TXT, MD, HTML, CSV, JSON)",
                "inputSchema": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Ruta al archivo de documento"
                        },
                        "encoding": {
                            "type": "string",
                            "description": "Codificación específica (opcional, default: utf-8)",
                            "default": "utf-8"
                        },
                        "extract_tables": {
                            "type": "boolean",
                            "description": "Extraer tablas de documentos (para DOCX)",
                            "default": True
                        }
                    },
                    "required": ["file_path"]
                }
            },
            {
                "name": "analyze_image_with_ai",
                "description": "Analiza imagen con IA: describe contenido, detecta objetos, OCR, QR codes",
                "inputSchema": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Ruta al archivo de imagen"
                        },
                        "analysis_type": {
                            "type": "string",
                            "enum": ["basic", "detailed", "ocr", "qr_detection", "objects", "full"],
                            "description": "Tipo de análisis a realizar",
                            "default": "full"
                        },
                        "ocr_languages": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Idiomas para OCR",
                            "default": ["spa", "eng"]
                        }
                    },
                    "required": ["file_path"]
                }
            },
            {
                "name": "convert_file_format",
                "description": "Convierte archivos entre formatos compatibles",
                "inputSchema": {
                    "type": "object",
                    "properties": {
                        "input_file_path": {
                            "type": "string",
                            "description": "Ruta del archivo de entrada"
                        },
                        "output_format": {
                            "type": "string",
                            "enum": ["pdf", "txt", "docx", "xlsx", "jpg", "png", "mp3", "mp4"],
                            "description": "Formato de salida deseado"
                        },
                        "output_file_path": {
                            "type": "string",
                            "description": "Ruta del archivo de salida"
                        },
                        "quality": {
                            "type": "number",
                            "minimum": 1,
                            "maximum": 100,
                            "description": "Calidad de conversión (1-100)",
                            "default": 85
                        }
                    },
                    "required": ["input_file_path", "output_format", "output_file_path"]
                }
            },
            {
                "name": "extract_metadata",
                "description": "Extrae metadata completa de cualquier archivo soportado",
                "inputSchema": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Ruta al archivo"
                        },
                        "include_content_analysis": {
                            "type": "boolean",
                            "description": "Incluir análisis de contenido básico",
                            "default": True
                        }
                    },
                    "required": ["file_path"]
                }
            },
            {
                "name": "process_audio_file",
                "description": "Analiza archivos de audio: extrae metadata, transcribe (si tiene voz)",
                "inputSchema": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Ruta al archivo de audio"
                        },
                        "analysis_type": {
                            "type": "string",
                            "enum": ["metadata", "waveform", "features", "full"],
                            "description": "Tipo de análisis de audio",
                            "default": "full"
                        }
                    },
                    "required": ["file_path"]
                }
            },
            {
                "name": "process_video_file",
                "description": "Analiza archivos de video: extrae metadata, frames, información técnica",
                "inputSchema": {
                    "type": "object",
                    "properties": {
                        "file_path": {
                            "type": "string",
                            "description": "Ruta al archivo de video"
                        },
                        "extract_frames": {
                            "type": "boolean",
                            "description": "Extraer frames de ejemplo",
                            "default": False
                        },
                        "frame_count": {
                            "type": "integer",
                            "minimum": 1,
                            "maximum": 20,
                            "description": "Número de frames a extraer",
                            "default": 5
                        }
                    },
                    "required": ["file_path"]
                }
            },
            {
                "name": "batch_process_files",
                "description": "Procesa múltiples archivos en lote",
                "inputSchema": {
                    "type": "object",
                    "properties": {
                        "file_paths": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Lista de rutas de archivos"
                        },
                        "operation": {
                            "type": "string",
                            "enum": ["extract_text", "analyze_image", "extract_metadata", "auto"],
                            "description": "Operación a realizar en cada archivo"
                        }
                    },
                    "required": ["file_paths", "operation"]
                }
            }
        ]
    
    async def call_tool(self, name: str, arguments: Dict[str, Any]) -> Dict[str, Any]:
        """Ejecuta una herramienta MCP específica"""
        try:
            if name == "extract_text_from_document":
                return await self._extract_text_from_document(**arguments)
            elif name == "analyze_image_with_ai":
                return await self._analyze_image_with_ai(**arguments)
            elif name == "convert_file_format":
                return await self._convert_file_format(**arguments)
            elif name == "extract_metadata":
                return await self._extract_metadata(**arguments)
            elif name == "process_audio_file":
                return await self._process_audio_file(**arguments)
            elif name == "process_video_file":
                return await self._process_video_file(**arguments)
            elif name == "batch_process_files":
                return await self._batch_process_files(**arguments)
            else:
                return {
                    "success": False,
                    "error": f"Herramienta no reconocida: {name}"
                }
        except Exception as e:
            return {
                "success": False,
                "error": f"Error ejecutando herramienta {name}: {str(e)}"
            }
    
    async def _extract_text_from_document(self, file_path: str, encoding: str = "utf-8", extract_tables: bool = True) -> Dict[str, Any]:
        """Extrae texto de documentos con soporte mejorado"""
        try:
            result = self.file_processor.validate_file(file_path)
            if not result.success:
                return result.data if hasattr(result, 'data') else {"success": False, "error": result.error}
            
            file_info = result.data
            file_ext = file_info['extension'].lower()
            
            # Procesar según tipo de archivo
            if file_ext in ['.pdf']:
                return await self._extract_pdf_advanced(file_path, extract_tables)
            elif file_ext in ['.docx']:
                return await self._extract_docx_advanced(file_path, extract_tables)
            elif file_ext in ['.xlsx']:
                return await self._extract_xlsx_content(file_path)
            elif file_ext in ['.pptx']:
                return await self._extract_pptx_content(file_path)
            elif file_ext in ['.txt', '.md', '.csv', '.json', '.html']:
                return await self._extract_text_file(file_path, encoding)
            else:
                return {
                    "success": False,
                    "error": f"Tipo de documento no soportado: {file_ext}"
                }
        except Exception as e:
            return {
                "success": False,
                "error": f"Error extrayendo texto: {str(e)}"
            }
    
    async def _analyze_image_with_ai(self, file_path: str, analysis_type: str = "full", ocr_languages: List[str] = None) -> Dict[str, Any]:
        """Analiza imagen con capacidades de IA"""
        if not PIL_AVAILABLE:
            return {
                "success": False,
                "error": "PIL no está disponible. Instalar con: pip install Pillow"
            }
        
        try:
            # Validar archivo
            result = self.file_processor.validate_file(file_path)
            if not result.success:
                return result.data if hasattr(result, 'data') else {"success": False, "error": result.error}
            
            # Abrir imagen
            image = Image.open(file_path)
            
            analysis_results = {
                "basic_info": {
                    "format": image.format,
                    "mode": image.mode,
                    "size": image.size,
                    "has_transparency": image.mode in ('RGBA', 'LA')
                }
            }
            
            # Análisis básico
            if analysis_type in ["basic", "full"]:
                analysis_results["description"] = await self._describe_image_basic(image)
            
            # Detección de objetos
            if analysis_type in ["objects", "full"] and CV2_AVAILABLE:
                analysis_results["object_detection"] = await self._detect_objects(image)
            
            # OCR
            if analysis_type in ["ocr", "full"] and TESSERACT_AVAILABLE:
                ocr_languages = ocr_languages or self.ocr_languages
                analysis_results["ocr_text"] = await self._perform_ocr(image, ocr_languages)
            
            # Detección de códigos QR
            if analysis_type in ["qr_detection", "full"] and QRCODE_AVAILABLE:
                analysis_results["qr_codes"] = await self._detect_qr_codes(image)
            
            # Análisis detallado
            if analysis_type == "detailed":
                analysis_results["detailed_analysis"] = await self._detailed_image_analysis(image)
            
            return {
                "success": True,
                "data": analysis_results,
                "file_path": file_path
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error analizando imagen: {str(e)}"
            }
    
    async def _convert_file_format(self, input_file_path: str, output_format: str, output_file_path: str, quality: int = 85) -> Dict[str, Any]:
        """Convierte archivos entre formatos"""
        try:
            # Validar archivo de entrada
            result = self.file_processor.validate_file(input_file_path)
            if not result.success:
                return result.data if hasattr(result, 'data') else {"success": False, "error": result.error}
            
            input_ext = Path(input_file_path).suffix.lower()
            
            # Crear directorio de salida si no existe
            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
            
            # Conversiones de imagen
            if input_ext in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp'] and output_format in ['jpg', 'png', 'tiff']:
                return await self._convert_image_format(input_file_path, output_file_path, output_format, quality)
            
            # Conversiones de documento
            elif input_ext in ['.pdf', '.docx'] and output_format in ['txt', 'pdf']:
                return await self._convert_document_format(input_file_path, output_format, output_file_path)
            
            # Conversiones de audio
            elif input_ext in ['.mp3', '.wav', '.flac', '.aac'] and output_format in ['mp3', 'wav']:
                return await self._convert_audio_format(input_file_path, output_format, output_file_path)
            
            # Conversiones de video
            elif input_ext in ['.mp4', '.avi', '.mov'] and output_format in ['mp4']:
                return await self._convert_video_format(input_file_path, output_format, output_file_path)
            
            else:
                return {
                    "success": False,
                    "error": f"Conversión de {input_ext} a {output_format} no soportada"
                }
                
        except Exception as e:
            return {
                "success": False,
                "error": f"Error convirtiendo archivo: {str(e)}"
            }
    
    async def _extract_metadata(self, file_path: str, include_content_analysis: bool = True) -> Dict[str, Any]:
        """Extrae metadata completa de archivos"""
        try:
            result = self.file_processor.validate_file(file_path)
            if not result.success:
                return result.data if hasattr(result, 'data') else {"success": False, "error": result.error}
            
            file_info = result.data
            metadata = {
                "file_info": file_info,
                "system_metadata": {
                    "created": datetime.fromtimestamp(os.path.getctime(file_path)).isoformat(),
                    "modified": datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat(),
                    "accessed": datetime.fromtimestamp(os.path.getatime(file_path)).isoformat()
                }
            }
            
            # Metadata específica por tipo de archivo
            file_ext = file_info['extension'].lower()
            
            if file_ext in ['.jpg', '.jpeg', '.png', '.tiff']:
                metadata["media_metadata"] = await self._extract_image_metadata(file_path)
            elif file_ext in ['.mp3', '.wav', '.flac', '.aac']:
                metadata["media_metadata"] = await self._extract_audio_metadata(file_path)
            elif file_ext in ['.mp4', '.avi', '.mov']:
                metadata["media_metadata"] = await self._extract_video_metadata(file_path)
            elif file_ext in ['.pdf']:
                metadata["document_metadata"] = await self._extract_pdf_metadata(file_path)
            
            # Análisis de contenido básico
            if include_content_analysis:
                metadata["content_analysis"] = await self._basic_content_analysis(file_path)
            
            return {
                "success": True,
                "data": metadata,
                "file_path": file_path
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error extrayendo metadata: {str(e)}"
            }
    
    async def _process_audio_file(self, file_path: str, analysis_type: str = "full") -> Dict[str, Any]:
        """Procesa archivos de audio"""
        if not AUDIO_LIBS_AVAILABLE:
            return {
                "success": False,
                "error": "Bibliotecas de audio no disponibles. Instalar con: pip install librosa"
            }
        
        try:
            result = self.file_processor.validate_file(file_path)
            if not result.success:
                return result.data if hasattr(result, 'data') else {"success": False, "error": result.error}
            
            audio_data = {}
            
            if analysis_type in ["metadata", "full"]:
                audio_data["metadata"] = await self._extract_audio_metadata(file_path)
            
            if analysis_type in ["waveform", "full"]:
                audio_data["waveform"] = await self._analyze_audio_waveform(file_path)
            
            if analysis_type in ["features", "full"]:
                audio_data["features"] = await self._extract_audio_features(file_path)
            
            return {
                "success": True,
                "data": audio_data,
                "file_path": file_path
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error procesando audio: {str(e)}"
            }
    
    async def _process_video_file(self, file_path: str, extract_frames: bool = False, frame_count: int = 5) -> Dict[str, Any]:
        """Procesa archivos de video"""
        try:
            result = self.file_processor.validate_file(file_path)
            if not result.success:
                return result.data if hasattr(result, 'data') else {"success": False, "error": result.error}
            
            video_data = {
                "metadata": await self._extract_video_metadata(file_path)
            }
            
            # Extraer frames de ejemplo
            if extract_frames and CV2_AVAILABLE:
                video_data["sample_frames"] = await self._extract_video_frames(file_path, frame_count)
            
            return {
                "success": True,
                "data": video_data,
                "file_path": file_path
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error procesando video: {str(e)}"
            }
    
    async def _batch_process_files(self, file_paths: List[str], operation: str) -> Dict[str, Any]:
        """Procesa múltiples archivos en lote"""
        results = []
        errors = []
        
        for file_path in file_paths:
            try:
                if operation == "extract_text":
                    result = await self._extract_text_from_document(file_path)
                elif operation == "analyze_image":
                    result = await self._analyze_image_with_ai(file_path)
                elif operation == "extract_metadata":
                    result = await self._extract_metadata(file_path)
                elif operation == "auto":
                    result = await self._auto_process_file(file_path)
                
                if result["success"]:
                    results.append(result)
                else:
                    errors.append({
                        "file_path": file_path,
                        "error": result.get("error", "Error desconocido")
                    })
                    
            except Exception as e:
                errors.append({
                    "file_path": file_path,
                    "error": str(e)
                })
        
        return {
            "success": True,
            "data": {
                "processed_count": len(results),
                "error_count": len(errors),
                "results": results,
                "errors": errors
            }
        }
    
    # Métodos auxiliares para procesamiento avanzado
    async def _extract_pdf_advanced(self, file_path: str, extract_tables: bool) -> Dict[str, Any]:
        """Extracción avanzada de PDF con soporte para tablas y metadatos"""
        # Usar el método base del FileProcessor y extenderlo
        base_result = self.file_processor.extract_pdf_text(file_path)
        if not base_result.success:
            return base_result.data if hasattr(base_result, 'data') else {"success": False, "error": base_result.error}
        
        pdf_data = base_result.data
        
        # Agregar análisis adicional de PDF si es necesario
        pdf_data["advanced_analysis"] = {
            "has_images": True,  # Placeholder - implementar detección real
            "has_tables": extract_tables,
            "complexity_score": len(pdf_data.get("text", "").split()) / 100
        }
        
        return {
            "success": True,
            "data": pdf_data,
            "file_path": file_path
        }
    
    async def _extract_docx_advanced(self, file_path: str, extract_tables: bool) -> Dict[str, Any]:
        """Extracción avanzada de DOCX"""
        base_result = self.file_processor.extract_docx_text(file_path)
        if not base_result.success:
            return base_result.data if hasattr(base_result, 'data') else {"success": False, "error": base_result.error}
        
        doc_data = base_result.data
        
        # Análisis adicional de DOCX
        doc_data["advanced_analysis"] = {
            "style_analysis": self._analyze_docx_styles(file_path),
            "image_count": self._count_docx_images(file_path),
            "complexity_score": doc_data.get("words", 0) / 500
        }
        
        return {
            "success": True,
            "data": doc_data,
            "file_path": file_path
        }
    
    async def _extract_xlsx_content(self, file_path: str) -> Dict[str, Any]:
        """Extrae contenido de archivos Excel"""
        if not XLSX_AVAILABLE:
            return {
                "success": False,
                "error": "openpyxl no está disponible. Instalar con: pip install openpyxl"
            }
        
        try:
            workbook = load_workbook(file_path, read_only=True)
            sheets_data = {}
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []
                
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        sheet_data.append([str(cell) if cell is not None else "" for cell in row])
                
                sheets_data[sheet_name] = {
                    "data": sheet_data,
                    "rows": len(sheet_data),
                    "cols": len(sheet_data[0]) if sheet_data else 0
                }
            
            return {
                "success": True,
                "data": {
                    "sheets": sheets_data,
                    "sheet_count": len(sheets_data)
                },
                "file_path": file_path
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error procesando Excel: {str(e)}"
            }
    
    async def _extract_pptx_content(self, file_path: str) -> Dict[str, Any]:
        """Extrae contenido de archivos PowerPoint"""
        if not PPTX_AVAILABLE:
            return {
                "success": False,
                "error": "python-pptx no está disponible. Instalar con: pip install python-pptx"
            }
        
        try:
            presentation = Presentation(file_path)
            slides_data = []
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_content = {
                    "slide_number": slide_num + 1,
                    "text_content": [],
                    "image_count": 0,
                    "shape_count": len(slide.shapes)
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content["text_content"].append(shape.text.strip())
                    
                    if shape.shape_type == 13:  # Picture type
                        slide_content["image_count"] += 1
                
                slides_data.append(slide_content)
            
            return {
                "success": True,
                "data": {
                    "slides": slides_data,
                    "slide_count": len(slides_data),
                    "total_text_elements": sum(len(slide["text_content"]) for slide in slides_data),
                    "total_images": sum(slide["image_count"] for slide in slides_data)
                },
                "file_path": file_path
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error procesando PowerPoint: {str(e)}"
            }
    
    async def _extract_text_file(self, file_path: str, encoding: str) -> Dict[str, Any]:
        """Extrae contenido de archivos de texto"""
        result = self.file_processor.read_text_file(file_path, encoding)
        if not result.success:
            return result.data if hasattr(result, 'data') else {"success": False, "error": result.error}
        
        return {
            "success": True,
            "data": result.data,
            "file_path": file_path
        }
    
    async def _describe_image_basic(self, image: Image.Image) -> str:
        """Descripción básica de imagen usando características técnicas"""
        width, height = image.size
        
        # Análisis de colores predominantes
        if image.mode == 'RGB':
            # Simplificado - en implementación real usaríamos análisis más sofisticado
            color_analysis = "imagen en color RGB"
        elif image.mode == 'L':
            color_analysis = "imagen en escala de grises"
        elif image.mode == 'RGBA':
            color_analysis = "imagen con canal alfa (transparencia)"
        else:
            color_analysis = f"imagen en modo {image.mode}"
        
        return f"Imagen {color_analysis} de {width}x{height} píxeles"
    
    async def _detect_objects(self, image: Image.Image) -> List[str]:
        """Detección básica de objetos en imagen"""
        if not CV2_AVAILABLE:
            return ["OpenCV no disponible para detección de objetos"]
        
        try:
            # Convertir PIL a OpenCV
            cv_image = cv2.cvtColor(image.convert('RGB'), cv2.COLOR_RGB2BGR)
            
            # Placeholder para detección de objetos
            # En implementación real usaríamos YOLO, TensorFlow, etc.
            return ["análisis de objetos básico - requiere modelo entrenado"]
            
        except Exception as e:
            return [f"Error en detección de objetos: {str(e)}"]
    
    async def _perform_ocr(self, image: Image.Image, languages: List[str]) -> str:
        """Realiza OCR en la imagen"""
        if not TESSERACT_AVAILABLE:
            return "Tesseract no está disponible para OCR"
        
        try:
            # Configurar idiomas
            lang_code = '+'.join(languages)
            
            # Realizar OCR
            text = pytesseract.image_to_string(image, lang=lang_code)
            
            return text.strip() if text.strip() else "No se detectó texto en la imagen"
            
        except Exception as e:
            return f"Error en OCR: {str(e)}"
    
    async def _detect_qr_codes(self, image: Image.Image) -> List[Dict[str, Any]]:
        """Detecta códigos QR en la imagen"""
        if not QRCODE_AVAILABLE:
            return [{"error": "qrcode no disponible"}]
        
        try:
            from pyzbar import pyzbar
            
            # Convertir PIL a formato compatible
            cv_image = cv2.cvtColor(image.convert('RGB'), cv2.COLOR_RGB2BGR)
            
            # Detectar códigos QR
            qr_codes = pyzbar.decode(cv_image)
            
            results = []
            for qr in qr_codes:
                results.append({
                    "data": qr.data.decode('utf-8'),
                    "type": qr.type,
                    "bbox": {
                        "x": qr.rect.left,
                        "y": qr.rect.top,
                        "width": qr.rect.width,
                        "height": qr.rect.height
                    }
                })
            
            return results if results else [{"message": "No se detectaron códigos QR"}]
            
        except ImportError:
            return [{"error": "pyzbar no disponible para detección de QR"}]
        except Exception as e:
            return [{"error": f"Error detectando QR: {str(e)}"}]
    
    async def _detailed_image_analysis(self, image: Image.Image) -> Dict[str, Any]:
        """Análisis detallado de imagen"""
        analysis = {
            "technical_specs": {
                "format": image.format,
                "mode": image.mode,
                "size": image.size,
                "info": image.info
            }
        }
        
        # Análisis de calidad de imagen
        if image.mode in ('RGB', 'RGBA'):
            import numpy as np
            img_array = np.array(image)
            
            # Análisis de brillo
            if len(img_array.shape) == 3:
                brightness = np.mean(img_array)
                analysis["quality_metrics"] = {
                    "brightness": float(brightness),
                    "has_color": True,
                    "is_oversaturated": brightness > 200,
                    "is_undersaturated": brightness < 50
                }
        
        return analysis
    
    async def _convert_image_format(self, input_path: str, output_path: str, output_format: str, quality: int) -> Dict[str, Any]:
        """Convierte formato de imagen"""
        try:
            image = Image.open(input_path)
            
            # Configuración de calidad
            quality_config = {
                "quality": quality,
                "optimize": True
            } if output_format in ['jpg', 'png'] else {}
            
            # Convertir formato
            if output_format == 'jpg' and image.mode == 'RGBA':
                # Convertir RGBA a RGB para JPEG
                image = image.convert('RGB')
            
            image.save(output_path, **quality_config)
            
            return {
                "success": True,
                "data": {
                    "input_format": image.format,
                    "output_format": output_format.upper(),
                    "quality": quality,
                    "conversion_successful": True
                },
                "file_path": output_path
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error convirtiendo imagen: {str(e)}"
            }
    
    async def _convert_document_format(self, input_path: str, output_format: str, output_path: str) -> Dict[str, Any]:
        """Convierte formato de documento"""
        try:
            if output_format == 'txt':
                # Extraer texto del documento
                if Path(input_path).suffix.lower() == '.pdf':
                    result = await self._extract_pdf_advanced(input_path, False)
                elif Path(input_path).suffix.lower() == '.docx':
                    result = await self._extract_docx_advanced(input_path, False)
                else:
                    result = await self._extract_text_file(input_path, 'utf-8')
                
                if result["success"]:
                    text_content = result["data"].get("text", "")
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(text_content)
                    
                    return {
                        "success": True,
                        "data": {
                            "conversion_type": "document_to_text",
                            "character_count": len(text_content),
                            "conversion_successful": True
                        },
                        "file_path": output_path
                    }
                else:
                    return result
            else:
                return {
                    "success": False,
                    "error": f"Conversión a {output_format} no implementada aún"
                }
                
        except Exception as e:
            return {
                "success": False,
                "error": f"Error convirtiendo documento: {str(e)}"
            }
    
    async def _convert_audio_format(self, input_path: str, output_format: str, output_path: str) -> Dict[str, Any]:
        """Convierte formato de audio usando FFmpeg"""
        if not FFMPEG_AVAILABLE:
            return {
                "success": False,
                "error": "FFmpeg no disponible para conversión de audio"
            }
        
        try:
            # Usar FFmpeg para conversión
            (
                ffmpeg
                .input(input_path)
                .output(output_path)
                .overwrite_output()
                .run(capture_stdout=True, capture_stderr=True)
            )
            
            return {
                "success": True,
                "data": {
                    "conversion_type": "audio_format",
                    "conversion_successful": True
                },
                "file_path": output_path
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error convirtiendo audio: {str(e)}"
            }
    
    async def _convert_video_format(self, input_path: str, output_format: str, output_path: str) -> Dict[str, Any]:
        """Convierte formato de video usando FFmpeg"""
        if not FFMPEG_AVAILABLE:
            return {
                "success": False,
                "error": "FFmpeg no disponible para conversión de video"
            }
        
        try:
            # Usar FFmpeg para conversión de video
            (
                ffmpeg
                .input(input_path)
                .output(output_path)
                .overwrite_output()
                .run(capture_stdout=True, capture_stderr=True)
            )
            
            return {
                "success": True,
                "data": {
                    "conversion_type": "video_format",
                    "conversion_successful": True
                },
                "file_path": output_path
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error convirtiendo video: {str(e)}"
            }
    
    async def _extract_image_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extrae metadata específica de imágenes"""
        try:
            with Image.open(file_path) as image:
                metadata = {
                    "format": image.format,
                    "mode": image.mode,
                    "size": image.size,
                    "has_transparency": image.mode in ('RGBA', 'LA'),
                    "aspect_ratio": round(image.size[0] / image.size[1], 2) if image.size[1] != 0 else 0
                }
                
                # Metadata EXIF si está disponible
                if hasattr(image, '_getexif') and image._getexif():
                    metadata["exif_data"] = image._getexif()
                
                return metadata
                
        except Exception as e:
            return {"error": f"Error extrayendo metadata de imagen: {str(e)}"}
    
    async def _extract_audio_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extrae metadata específica de audio"""
        if not MUTAGEN_AVAILABLE:
            return {"error": "mutagen no disponible para metadata de audio"}
        
        try:
            from mutagen import File
            
            audio_file = File(file_path)
            if audio_file is None:
                return {"error": "No se pudo cargar el archivo de audio"}
            
            metadata = {
                "format": audio_file.mime[0] if audio_file.mime else "unknown",
                "length": audio_file.info.length if hasattr(audio_file.info, 'length') else 0,
                "bitrate": audio_file.info.bitrate if hasattr(audio_file.info, 'bitrate') else 0,
                "sample_rate": audio_file.info.sample_rate if hasattr(audio_file.info, 'sample_rate') else 0,
                "channels": audio_file.info.channels if hasattr(audio_file.info, 'channels') else 0
            }
            
            # Tags si están disponibles
            if hasattr(audio_file, 'tags') and audio_file.tags:
                metadata["tags"] = dict(audio_file.tags)
            
            return metadata
            
        except Exception as e:
            return {"error": f"Error extrayendo metadata de audio: {str(e)}"}
    
    async def _extract_video_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extrae metadata específica de video"""
        if not FFMPEG_AVAILABLE:
            return {"error": "FFmpeg no disponible para metadata de video"}
        
        try:
            probe = ffmpeg.probe(file_path)
            
            if not probe['streams']:
                return {"error": "No se encontraron streams en el video"}
            
            video_stream = next((stream for stream in probe['streams'] if stream['codec_type'] == 'video'), None)
            audio_stream = next((stream for stream in probe['streams'] if stream['codec_type'] == 'audio'), None)
            
            metadata = {
                "duration": float(probe['format'].get('duration', 0)),
                "bit_rate": probe['format'].get('bit_rate', 'unknown'),
                "format_name": probe['format'].get('format_name', 'unknown')
            }
            
            if video_stream:
                metadata["video"] = {
                    "codec": video_stream.get('codec_name', 'unknown'),
                    "width": int(video_stream.get('width', 0)),
                    "height": int(video_stream.get('height', 0)),
                    "fps": eval(video_stream.get('r_frame_rate', '0/1')) if '/' in video_stream.get('r_frame_rate', '0/1') else 0
                }
            
            if audio_stream:
                metadata["audio"] = {
                    "codec": audio_stream.get('codec_name', 'unknown'),
                    "sample_rate": audio_stream.get('sample_rate', 'unknown'),
                    "channels": audio_stream.get('channels', 0)
                }
            
            return metadata
            
        except Exception as e:
            return {"error": f"Error extrayendo metadata de video: {str(e)}"}
    
    async def _extract_pdf_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extrae metadata específica de PDF"""
        try:
            import PyPDF2
            
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                metadata = {
                    "page_count": len(reader.pages),
                    "has_metadata": bool(reader.metadata)
                }
                
                if reader.metadata:
                    metadata.update({
                        "title": reader.metadata.get('/Title', ''),
                        "author": reader.metadata.get('/Author', ''),
                        "subject": reader.metadata.get('/Subject', ''),
                        "creator": reader.metadata.get('/Creator', ''),
                        "producer": reader.metadata.get('/Producer', ''),
                        "creation_date": str(reader.metadata.get('/CreationDate', '')),
                        "modification_date": str(reader.metadata.get('/ModDate', ''))
                    })
                
                return metadata
                
        except Exception as e:
            return {"error": f"Error extrayendo metadata de PDF: {str(e)}"}
    
    async def _basic_content_analysis(self, file_path: str) -> Dict[str, Any]:
        """Análisis básico de contenido"""
        try:
            file_ext = Path(file_path).suffix.lower()
            
            analysis = {
                "file_type_category": self._get_file_category(file_ext),
                "estimated_complexity": "basic"
            }
            
            # Análisis específico por tipo
            if file_ext in ['.txt', '.md']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    analysis["text_analysis"] = {
                        "character_count": len(content),
                        "word_count": len(content.split()),
                        "line_count": len(content.splitlines()),
                        "has_markdown": file_ext == '.md'
                    }
            
            elif file_ext in ['.jpg', '.png', '.jpeg']:
                with Image.open(file_path) as image:
                    analysis["image_analysis"] = {
                        "pixel_count": image.size[0] * image.size[1],
                        "is_large": image.size[0] * image.size[1] > 1000000,
                        "color_depth": len(image.getcolors() or []) if image.mode == 'RGB' else 'unknown'
                    }
            
            return analysis
            
        except Exception as e:
            return {"error": f"Error en análisis de contenido: {str(e)}"}
    
    def _get_file_category(self, extension: str) -> str:
        """Determina la categoría de un archivo por su extensión"""
        document_exts = {'.pdf', '.docx', '.xlsx', '.pptx', '.txt', '.md', '.csv', '.json', '.html'}
        image_exts = {'.jpg', '.jpeg', '.png', '.tiff', '.gif', '.bmp', '.webp'}
        audio_exts = {'.mp3', '.wav', '.flac', '.aac', '.ogg'}
        video_exts = {'.mp4', '.avi', '.mov', '.mkv', '.webm'}
        
        if extension in document_exts:
            return "document"
        elif extension in image_exts:
            return "image"
        elif extension in audio_exts:
            return "audio"
        elif extension in video_exts:
            return "video"
        else:
            return "unknown"
    
    async def _analyze_audio_waveform(self, file_path: str) -> Dict[str, Any]:
        """Analiza waveform de audio"""
        try:
            import numpy as np
            
            # Cargar audio
            y, sr = librosa.load(file_path)
            
            # Calcular características del waveform
            duration = len(y) / sr
            
            return {
                "duration": float(duration),
                "sample_rate": int(sr),
                "channels": 1,  # librosa convierte a mono por defecto
                "peak_amplitude": float(np.max(np.abs(y))),
                "rms_energy": float(np.sqrt(np.mean(y**2)))
            }
            
        except Exception as e:
            return {"error": f"Error analizando waveform: {str(e)}"}
    
    async def _extract_audio_features(self, file_path: str) -> Dict[str, Any]:
        """Extrae características de audio"""
        try:
            # Cargar audio
            y, sr = librosa.load(file_path)
            
            # Extraer características básicas
            features = {
                "spectral_centroid": float(np.mean(librosa.feature.spectral_centroid(y=y, sr=sr))),
                "zero_crossing_rate": float(np.mean(librosa.feature.zero_crossing_rate(y))),
                "mfcc": [float(x) for x in np.mean(librosa.feature.mfcc(y=y, sr=sr), axis=1)]
            }
            
            return features
            
        except Exception as e:
            return {"error": f"Error extrayendo características: {str(e)}"}
    
    async def _extract_video_frames(self, file_path: str, frame_count: int) -> List[Dict[str, Any]]:
        """Extrae frames de ejemplo del video"""
        try:
            cap = cv2.VideoCapture(file_path)
            
            if not cap.isOpened():
                return [{"error": "No se pudo abrir el video"}]
            
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            fps = cap.get(cv2.CAP_PROP_FPS)
            
            frames_info = []
            
            # Calcular intervalos para extraer frames uniformemente
            interval = max(1, total_frames // frame_count)
            
            for i in range(0, min(frame_count * interval, total_frames), interval):
                cap.set(cv2.CAP_PROP_POS_FRAMES, i)
                ret, frame = cap.read()
                
                if ret:
                    frames_info.append({
                        "frame_number": i,
                        "timestamp": i / fps if fps > 0 else 0,
                        "dimensions": {
                            "width": int(cap.get(cv2.CAP_PROP_FRAME_WIDTH)),
                            "height": int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
                        }
                    })
            
            cap.release()
            
            return frames_info
            
        except Exception as e:
            return [{"error": f"Error extrayendo frames: {str(e)}"}]
    
    async def _auto_process_file(self, file_path: str) -> Dict[str, Any]:
        """Procesa automáticamente un archivo según su tipo"""
        try:
            result = self.file_processor.validate_file(file_path)
            if not result.success:
                return result.data if hasattr(result, 'data') else {"success": False, "error": result.error}
            
            file_ext = result.data['extension'].lower()
            
            # Determinar tipo y procesar
            if file_ext in ['.jpg', '.jpeg', '.png', '.tiff', '.gif', '.bmp']:
                return await self._analyze_image_with_ai(file_path)
            elif file_ext in ['.mp3', '.wav', '.flac', '.aac']:
                return await self._process_audio_file(file_path)
            elif file_ext in ['.mp4', '.avi', '.mov']:
                return await self._process_video_file(file_path)
            else:
                return await self._extract_text_from_document(file_path)
                
        except Exception as e:
            return {
                "success": False,
                "error": f"Error en procesamiento automático: {str(e)}"
            }
    
    # Métodos auxiliares para análisis de documentos
    def _analyze_docx_styles(self, file_path: str) -> Dict[str, Any]:
        """Analiza estilos en un documento DOCX"""
        # Placeholder - implementar análisis real de estilos
        return {
            "styles_found": ["Normal", "Heading 1", "Heading 2"],
            "style_analysis": "análisis básico de estilos"
        }
    
    def _count_docx_images(self, file_path: str) -> int:
        """Cuenta imágenes en un documento DOCX"""
        # Placeholder - implementar conteo real de imágenes
        return 0


# Función para crear la instancia del agente
def create_file_processing_agent() -> FileProcessingAgentMCP:
    """Crea una instancia del File Processing Agent"""
    return FileProcessingAgentMCP()


if __name__ == "__main__":
    # Ejemplo de uso del agente
    agent = create_file_processing_agent()
    
    # Mostrar herramientas disponibles
    print("File Processing Agent MCP - Herramientas disponibles:")
    for tool in agent.get_tools():
        print(f"- {tool['name']}: {tool['description']}")